#!/usr/bin/env python3
"""Capture and compare protected PPTX text/image geometry for final-polish work.

Requires python-pptx. It intentionally reports differences rather than deciding
whether they were allowed: authorization belongs in baseline-protection.md and
the revision ledger.
"""

from __future__ import annotations

import argparse
import hashlib
import json
import sys
from pathlib import Path


def _sha256(blob: bytes) -> str:
    return hashlib.sha256(blob).hexdigest()


def _shape_record(shape, index: int) -> dict:
    record = {
        "shape_index": index,
        "shape_type": str(shape.shape_type),
        "name": shape.name,
        "bbox": [int(shape.left), int(shape.top), int(shape.width), int(shape.height)],
    }
    if getattr(shape, "has_text_frame", False):
        text = shape.text_frame.text.strip()
        if text:
            record["text"] = text
    if str(shape.shape_type).endswith("PICTURE (13)"):
        image = shape.image
        record["image_sha256"] = _sha256(image.blob)
        record["crop"] = [
            float(shape.crop_left), float(shape.crop_top),
            float(shape.crop_right), float(shape.crop_bottom),
        ]
    return record


def capture(deck: Path) -> dict:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("python-pptx is required; install it in the PPTX runtime.") from exc
    presentation = Presentation(str(deck))
    return {
        "schema": 1,
        "source": str(deck),
        "slide_count": len(presentation.slides),
        "slides": [
            {"slide": number, "shapes": [_shape_record(shape, i) for i, shape in enumerate(slide.shapes, 1)]}
            for number, slide in enumerate(presentation.slides, 1)
        ],
    }


def _key(record: dict) -> tuple[int, int]:
    return record["slide"], record["shape_index"]


def compare(manifest: dict, revised: dict) -> list[dict]:
    baseline = {
        (slide["slide"], shape["shape_index"]): shape
        for slide in manifest["slides"] for shape in slide["shapes"]
        if "text" in shape or "image_sha256" in shape
    }
    current = {
        (slide["slide"], shape["shape_index"]): shape
        for slide in revised["slides"] for shape in slide["shapes"]
        if "text" in shape or "image_sha256" in shape
    }
    differences = []
    for key, before in baseline.items():
        after = current.get(key)
        label = f"slide {key[0]}, shape {key[1]} ({before['name']})"
        if after is None:
            differences.append({"object": label, "change": "missing", "before": before, "after": None})
            continue
        for field in ("bbox", "text", "image_sha256", "crop"):
            if before.get(field) != after.get(field):
                differences.append({"object": label, "change": field, "before": before.get(field), "after": after.get(field)})
    return differences


def markdown_report(differences: list[dict], baseline: dict, revised: dict) -> str:
    lines = [
        "# Protected-Element Comparison",
        "",
        f"- Baseline slides: {baseline['slide_count']}",
        f"- Revised slides: {revised['slide_count']}",
        f"- Differences: {len(differences)}",
        "",
        "> Review every difference against `baseline-protection.md` and the revision ledger. A listed difference is not automatically authorized.",
        "",
    ]
    if not differences:
        return "\n".join(lines + ["No protected text/image geometry differences detected.", ""])
    lines.extend(["| Object | Change | Before | After | Authorized? |", "|---|---|---|---|---|"])
    for item in differences:
        before = str(item["before"]).replace("|", "\\|")[:120]
        after = str(item["after"]).replace("|", "\\|")[:120]
        lines.append(f"| {item['object']} | {item['change']} | {before} | {after} |  |")
    return "\n".join(lines) + "\n"


def main() -> int:
    parser = argparse.ArgumentParser(description="Capture/compare protected PPTX elements.")
    sub = parser.add_subparsers(dest="command", required=True)
    capture_parser = sub.add_parser("capture", help="write a baseline manifest")
    capture_parser.add_argument("deck", type=Path)
    capture_parser.add_argument("manifest", type=Path)
    compare_parser = sub.add_parser("compare", help="compare a manifest with a revised PPTX")
    compare_parser.add_argument("manifest", type=Path)
    compare_parser.add_argument("deck", type=Path)
    compare_parser.add_argument("--report", type=Path, required=True)
    args = parser.parse_args()
    try:
        if args.command == "capture":
            result = capture(args.deck)
            args.manifest.write_text(json.dumps(result, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
            print(f"captured {result['slide_count']} slides → {args.manifest}")
            return 0
        baseline = json.loads(args.manifest.read_text(encoding="utf-8"))
        revised = capture(args.deck)
        differences = compare(baseline, revised)
        args.report.write_text(markdown_report(differences, baseline, revised), encoding="utf-8")
        print(f"reported {len(differences)} difference(s) → {args.report}")
        return 1 if differences else 0
    except (OSError, ValueError, RuntimeError) as exc:
        print(f"error: {exc}", file=sys.stderr)
        return 2


if __name__ == "__main__":
    raise SystemExit(main())
